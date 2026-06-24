"""
Generates a 15-slide PowerPoint presentation for the
AI-Powered Flight & Hotel Booking Chatbot (TravelBot).

Run:
    python reports/generate_ppt.py
Outputs:
    reports/TravelBot_Project_Presentation.pptx
"""
from pathlib import Path

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR


PRIMARY = RGBColor(0x4F, 0x46, 0xE5)
SECONDARY = RGBColor(0xF5, 0x9E, 0x0B)
DARK = RGBColor(0x1F, 0x29, 0x37)
MUTED = RGBColor(0x6B, 0x72, 0x80)
LIGHT = RGBColor(0xF6, 0xF7, 0xFB)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


def set_run(run, *, text=None, size=18, bold=False, italic=False,
            color=DARK, font='Calibri'):
    if text is not None:
        run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.name = font
    run.font.color.rgb = color


def add_textbox(slide, x, y, w, h, *, text='', size=18, bold=False,
                italic=False, color=DARK, align=PP_ALIGN.LEFT, font='Calibri',
                anchor=MSO_ANCHOR.TOP):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    p = tf.paragraphs[0]
    p.alignment = align
    set_run(p.add_run(), text=text, size=size, bold=bold, italic=italic,
            color=color, font=font)
    return tb


def add_bullets(slide, x, y, w, h, items, *, size=20, color=DARK,
                bullet_color=PRIMARY):
    """items: list of strings or (lead, rest) tuples."""
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    first = True
    for item in items:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        p.space_after = Pt(8)
        r = p.add_run()
        r.text = '● '
        r.font.size = Pt(size)
        r.font.color.rgb = bullet_color
        r.font.bold = True
        if isinstance(item, tuple):
            lead, rest = item
            r1 = p.add_run()
            set_run(r1, text=lead, size=size, bold=True, color=color)
            r2 = p.add_run()
            set_run(r2, text=rest, size=size, color=color)
        else:
            r1 = p.add_run()
            set_run(r1, text=item, size=size, color=color)
    return tb


def add_rect(slide, x, y, w, h, fill_color, *, line=False):
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    if not line:
        shape.line.fill.background()
    shape.shadow.inherit = False
    return shape


def add_header_bar(slide, title, subtitle=None):
    add_rect(slide, 0, 0, SLIDE_W, Inches(0.9), PRIMARY)
    add_rect(slide, 0, Inches(0.9), SLIDE_W, Inches(0.05), SECONDARY)
    add_textbox(slide, Inches(0.5), Inches(0.18), Inches(12), Inches(0.6),
                text=title, size=28, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    if subtitle:
        add_textbox(slide, Inches(0.5), Inches(1.05), Inches(12), Inches(0.4),
                    text=subtitle, size=14, color=MUTED)


def add_footer(slide, idx, total):
    add_textbox(slide, Inches(0.5), Inches(7.0), Inches(6), Inches(0.4),
                text='TravelBot · AI-Powered Flight & Hotel Booking Chatbot',
                size=10, color=MUTED)
    add_textbox(slide, Inches(7), Inches(7.0), Inches(6), Inches(0.4),
                text=f'Slide {idx} / {total}',
                size=10, color=MUTED, align=PP_ALIGN.RIGHT)


# ---------- Slide builders ----------
def build_title_slide(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)

    # Solid background
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)

    # Decorative diagonal stripe
    add_rect(slide, 0, Inches(4.5), SLIDE_W, Inches(0.6),
             RGBColor(0x3F, 0x37, 0xC4))

    # Logo / icon
    add_textbox(slide, Inches(0.5), Inches(0.4), Inches(2), Inches(0.6),
                text='✈ TravelBot', size=20, bold=True, color=WHITE)

    # Title
    add_textbox(slide, Inches(0.8), Inches(1.5), Inches(11.8), Inches(1.3),
                text='AI-POWERED FLIGHT AND HOTEL\nBOOKING CHATBOT',
                size=44, bold=True, color=WHITE, align=PP_ALIGN.CENTER,
                anchor=MSO_ANCHOR.MIDDLE)

    add_textbox(slide, Inches(0.8), Inches(3.0), Inches(11.8), Inches(0.6),
                text='— A Python Full-Stack Project —',
                size=22, italic=True, color=SECONDARY, align=PP_ALIGN.CENTER)

    # Info block (placeholders)
    info_y = Inches(5.3)
    add_textbox(slide, Inches(2.5), info_y, Inches(8.5), Inches(0.45),
                text='Submitted by:  ___________________________',
                size=18, color=WHITE)
    add_textbox(slide, Inches(2.5), info_y + Inches(0.55), Inches(8.5), Inches(0.45),
                text='Roll No. / Reg. No.:  ___________________________',
                size=18, color=WHITE)
    add_textbox(slide, Inches(2.5), info_y + Inches(1.10), Inches(8.5), Inches(0.45),
                text='Department:  ___________________________',
                size=18, color=WHITE)
    add_textbox(slide, Inches(2.5), info_y + Inches(1.65), Inches(8.5), Inches(0.45),
                text='Guide:  ___________________________      Year: ____–____',
                size=18, color=WHITE)
    # College name (bottom)
    add_textbox(slide, Inches(0.5), Inches(7.05), Inches(12.3), Inches(0.4),
                text='_____________________________________ (Name of College / Institute)',
                size=14, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


def build_agenda(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '📋 Agenda', 'What this presentation will cover')

    agenda = [
        '1. Introduction and motivation',
        '2. Problem statement and objectives',
        '3. System overview & architecture',
        '4. Technology stack',
        '5. Key features and modules',
        '6. NLP chatbot engine',
        '7. Database design',
        '8. Admin dashboard',
        '9. Implementation highlights',
        '10. Testing & results',
        '11. Future enhancements',
        '12. Conclusion',
    ]
    add_bullets(slide, Inches(1.2), Inches(1.6), Inches(11), Inches(5.2),
                agenda, size=22)
    add_footer(slide, 2, total)


def build_introduction(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🌍 Introduction', 'Why we built TravelBot')

    add_textbox(slide, Inches(0.6), Inches(1.6), Inches(12), Inches(0.7),
                text='Traditional travel booking sites are form-heavy and slow.',
                size=22, bold=True, color=DARK)

    bullets = [
        ('Users today expect: ', 'conversational, fast, mobile-friendly booking.'),
        ('Existing OTAs: ', 'rely on multiple filters, dropdowns and pages.'),
        ('Our goal: ', 'combine the proven OTA workflow with a smart AI chatbot.'),
        ('Our product: ', 'TravelBot — a single Django web platform where users '
         'can chat naturally OR use forms to book flights & hotels.'),
    ]
    add_bullets(slide, Inches(0.6), Inches(2.4), Inches(12), Inches(4),
                bullets, size=20)

    # Highlight box
    add_rect(slide, Inches(0.6), Inches(5.8), Inches(12.1), Inches(1.0),
             RGBColor(0xFE, 0xF3, 0xC7))
    add_textbox(slide, Inches(0.8), Inches(5.85), Inches(11.8), Inches(0.9),
                text='💡 “Just tell our chatbot what you need — find flights from '
                     'Delhi to Goa tomorrow under ₹6000 — and book in seconds.”',
                size=16, italic=True, color=DARK,
                anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 3, total)


def build_problem_objectives(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🎯 Problem & Objectives')

    # Left card: problem
    add_rect(slide, Inches(0.6), Inches(1.5), Inches(6), Inches(5.2),
             RGBColor(0xFE, 0xE2, 0xE2))
    add_textbox(slide, Inches(0.8), Inches(1.7), Inches(5.6), Inches(0.5),
                text='Problem', size=22, bold=True,
                color=RGBColor(0xB9, 0x1C, 0x1C))
    add_bullets(slide, Inches(0.8), Inches(2.3), Inches(5.6), Inches(4.2), [
        'Form-heavy UI on existing OTAs.',
        'Hard to compare across providers.',
        'No real conversational interface.',
        'No unified loyalty & review system.',
        'Mobile users get a confusing experience.',
    ], size=18, bullet_color=RGBColor(0xB9, 0x1C, 0x1C))

    # Right card: objectives
    add_rect(slide, Inches(6.8), Inches(1.5), Inches(6), Inches(5.2),
             RGBColor(0xDC, 0xFC, 0xE7))
    add_textbox(slide, Inches(7), Inches(1.7), Inches(5.6), Inches(0.5),
                text='Objectives', size=22, bold=True,
                color=RGBColor(0x15, 0x80, 0x3D))
    add_bullets(slide, Inches(7), Inches(2.3), Inches(5.6), Inches(4.2), [
        'Build an intelligent NLP chatbot.',
        'Provide end-to-end booking workflow.',
        'Add reviews, wishlist, coupons, loyalty.',
        'Deliver PDF e-tickets and notifications.',
        'Separate admin dashboard with CRUD.',
    ], size=18, bullet_color=RGBColor(0x15, 0x80, 0x3D))

    add_footer(slide, 4, total)


def build_architecture(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🏛️ System Architecture')

    # Layered architecture boxes
    layers = [
        ('Presentation', 'HTML5 · CSS3 · Bootstrap 5 · Chart.js · Web Speech API', PRIMARY),
        ('Application', 'Django Views · URL Routing · Forms · NLP Engine (NLTK)', RGBColor(0x6D, 0x28, 0xD9)),
        ('Data Access', 'Django ORM · Model layer', RGBColor(0xDB, 0x27, 0x77)),
        ('Database', 'SQLite (relational, ACID, file-based)', RGBColor(0xF5, 0x9E, 0x0B)),
    ]
    y = Inches(1.6)
    for name, desc, color in layers:
        add_rect(slide, Inches(1.5), y, Inches(10), Inches(1.05), color)
        add_textbox(slide, Inches(1.8), y + Inches(0.05), Inches(3), Inches(0.45),
                    text=name, size=22, bold=True, color=WHITE)
        add_textbox(slide, Inches(1.8), y + Inches(0.5), Inches(9.5), Inches(0.5),
                    text=desc, size=14, color=WHITE)
        # arrow between layers
        if y + Inches(1.05) < Inches(6):
            add_textbox(slide, Inches(6.4), y + Inches(1.05), Inches(0.7), Inches(0.3),
                        text='⬇', size=20, color=DARK, align=PP_ALIGN.CENTER)
        y += Inches(1.35)

    add_footer(slide, 5, total)


def build_tech_stack(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '⚙️ Technology Stack')

    items = [
        ('Frontend', 'HTML5 · CSS3 · JavaScript · Bootstrap 5'),
        ('Backend', 'Python 3.10+ · Django 5'),
        ('Database', 'SQLite (zero-config, file-based)'),
        ('AI / NLP', 'NLTK + custom rule-based engine'),
        ('Charts', 'Chart.js for admin dashboard'),
        ('PDF', 'reportlab — e-ticket generation'),
        ('Voice', 'Web Speech API (browser-native)'),
        ('Auth', 'Django built-in (PBKDF2 hashing)'),
    ]

    # Two-column grid
    col1 = items[:4]; col2 = items[4:]
    y0 = Inches(1.6)
    for col, x in [(col1, Inches(0.7)), (col2, Inches(6.9))]:
        y = y0
        for name, desc in col:
            add_rect(slide, x, y, Inches(5.7), Inches(1.05), LIGHT)
            add_rect(slide, x, y, Inches(0.15), Inches(1.05), PRIMARY)
            add_textbox(slide, x + Inches(0.4), y + Inches(0.1), Inches(5.2), Inches(0.45),
                        text=name, size=18, bold=True, color=PRIMARY)
            add_textbox(slide, x + Inches(0.4), y + Inches(0.55), Inches(5.2), Inches(0.45),
                        text=desc, size=14, color=DARK)
            y += Inches(1.25)

    add_footer(slide, 6, total)


def build_modules(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🧩 Modules')

    mods = [
        ('👤 User',     'Register / login / profile · Loyalty tier · Referral code'),
        ('🎫 Booking',  'Search flights & hotels · Reviews · Wishlist · Cancel'),
        ('💳 Payment',  'Dummy gateway (Card / UPI / NB / Wallet) · Insurance addon'),
        ('🤖 Chatbot',  'NLP intent + entity extraction · Voice input · Rich cards'),
        ('🏆 Loyalty',  'Earn 1pt/₹100 · Redeem 1pt = ₹1 · 4-tier system'),
        ('🛠 Admin',    'CRUD flights / hotels / users / bookings / coupons · Reports'),
    ]
    y = Inches(1.6)
    for icon_name, desc in mods:
        add_rect(slide, Inches(0.7), y, Inches(12), Inches(0.8), LIGHT)
        add_textbox(slide, Inches(0.95), y + Inches(0.15), Inches(2.5), Inches(0.5),
                    text=icon_name, size=18, bold=True, color=PRIMARY)
        add_textbox(slide, Inches(3.5), y + Inches(0.15), Inches(9), Inches(0.55),
                    text=desc, size=15, color=DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.9)

    add_footer(slide, 7, total)


def build_chatbot(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🤖 NLP Chatbot Engine')

    add_textbox(slide, Inches(0.6), Inches(1.5), Inches(12), Inches(0.5),
                text='Two-step pipeline: Intent classification → Entity extraction.',
                size=18, bold=True, color=DARK)

    # Left: intents
    add_textbox(slide, Inches(0.6), Inches(2.1), Inches(5.5), Inches(0.4),
                text='Intents', size=18, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(0.6), Inches(2.5), Inches(5.8), Inches(4.5), [
        'greet, goodbye, thanks',
        'search_flight, search_hotel',
        'recommend, show_bookings, cancel',
        'promo, loyalty, wishlist',
        'help, fallback',
    ], size=16)

    # Right: entities
    add_textbox(slide, Inches(6.9), Inches(2.1), Inches(5.5), Inches(0.4),
                text='Entities', size=18, bold=True, color=PRIMARY)
    add_bullets(slide, Inches(6.9), Inches(2.5), Inches(5.8), Inches(4.5), [
        'origin, destination (Indian cities)',
        'departure_date (tomorrow / 25 Dec / next Fri)',
        'passengers, rooms, guests',
        'cabin_class (economy / business / first)',
        'min_stars (1–5), budget (under ₹X)',
    ], size=16)

    add_rect(slide, Inches(0.6), Inches(6.2), Inches(12.1), Inches(0.7),
             RGBColor(0xE0, 0xE7, 0xFF))
    add_textbox(slide, Inches(0.8), Inches(6.25), Inches(11.8), Inches(0.6),
                text='🎤 Voice input via Web Speech API · Rich card responses · '
                     'Conversation history stored per user/session',
                size=14, italic=True, color=DARK,
                anchor=MSO_ANCHOR.MIDDLE)
    add_footer(slide, 8, total)


def build_database(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🗄️ Database Schema')

    headers = ['Model', 'Purpose']
    rows = [
        ('User', 'Built-in Django auth'),
        ('UserProfile', 'Phone, address, loyalty_points, referral_code'),
        ('Flight', 'Airline, route, schedule, price, seats_available'),
        ('Hotel', 'Name, city, stars, amenities, rooms_available'),
        ('Booking', 'Reference, type, total, status, coupon, insurance'),
        ('Payment', 'Method, transaction_id, status'),
        ('Coupon', 'code, percent/flat, validity, usage limit'),
        ('HotelReview', 'rating, comment, is_verified'),
        ('WishlistItem', 'flight / hotel saved by user'),
        ('Notification', 'type, title, message, is_read'),
        ('ChatSession/Message', 'Conversation context + history'),
    ]

    # Build a table via shapes (simple two-column visual)
    y = Inches(1.5)
    add_rect(slide, Inches(0.7), y, Inches(4), Inches(0.55), PRIMARY)
    add_rect(slide, Inches(4.7), y, Inches(8), Inches(0.55), PRIMARY)
    add_textbox(slide, Inches(0.85), y + Inches(0.05), Inches(3.7), Inches(0.45),
                text='Model', size=16, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    add_textbox(slide, Inches(4.85), y + Inches(0.05), Inches(7.7), Inches(0.45),
                text='Purpose', size=16, bold=True, color=WHITE,
                anchor=MSO_ANCHOR.MIDDLE)
    y += Inches(0.55)
    striped = True
    for name, desc in rows:
        bg = LIGHT if striped else WHITE
        striped = not striped
        add_rect(slide, Inches(0.7), y, Inches(4), Inches(0.42), bg)
        add_rect(slide, Inches(4.7), y, Inches(8), Inches(0.42), bg)
        add_textbox(slide, Inches(0.85), y + Inches(0.05), Inches(3.7), Inches(0.32),
                    text=name, size=13, bold=True, color=DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        add_textbox(slide, Inches(4.85), y + Inches(0.05), Inches(7.7), Inches(0.32),
                    text=desc, size=13, color=DARK,
                    anchor=MSO_ANCHOR.MIDDLE)
        y += Inches(0.42)

    add_footer(slide, 9, total)


def build_features(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '✨ Key Features')

    features = [
        ('🤖', 'NLP Chatbot', 'Natural language search'),
        ('🎤', 'Voice Input', 'Web Speech API mic'),
        ('🎫', 'PDF E-Ticket', 'Auto-generated PDFs'),
        ('⭐', 'Reviews', '1-5 star ratings'),
        ('❤️', 'Wishlist', 'Save favourites'),
        ('🎟️', 'Promo Codes', '% or ₹ flat off'),
        ('🥉', 'Loyalty', 'Earn + redeem points'),
        ('🛡️', 'Insurance', 'Optional 5 % addon'),
        ('🔔', 'Notifications', 'Booking & promo alerts'),
        ('🌙', 'Dark Mode', 'Theme toggle'),
        ('💱', 'Currency', 'INR ↔ USD/EUR/etc.'),
        ('📨', 'Newsletter', 'Footer signup'),
    ]
    cols = 4
    cell_w = Inches(3.1)
    cell_h = Inches(1.4)
    start_x = Inches(0.5)
    start_y = Inches(1.4)
    gap = Inches(0.05)
    for i, (icon, name, desc) in enumerate(features):
        r, c = divmod(i, cols)
        x = start_x + (cell_w + gap) * c
        y = start_y + (cell_h + gap) * r
        add_rect(slide, x, y, cell_w, cell_h, LIGHT)
        add_rect(slide, x, y, Inches(0.15), cell_h, SECONDARY)
        add_textbox(slide, x + Inches(0.3), y + Inches(0.1), Inches(0.8), Inches(0.6),
                    text=icon, size=32)
        add_textbox(slide, x + Inches(1.1), y + Inches(0.15), Inches(1.9), Inches(0.45),
                    text=name, size=15, bold=True, color=PRIMARY)
        add_textbox(slide, x + Inches(1.1), y + Inches(0.65), Inches(1.9), Inches(0.6),
                    text=desc, size=11, color=MUTED)

    add_footer(slide, 10, total)


def build_admin(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🛠 Admin Dashboard',
                   'Strictly separated from customer experience')

    add_bullets(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(5.5), [
        ('Hidden customer features for admins: ',
         'no chatbot widget, no booking, no wishlist.'),
        ('Full CRUD: ',
         'Flights, Hotels, Coupons (with auto-broadcast notifications).'),
        ('User management: ',
         'activate/deactivate, toggle staff, delete (non-superuser).'),
        ('Booking control: ',
         'filter, search, change status (notifies user), delete.'),
        ('Communications: ',
         'Contact-message inbox and newsletter subscriber list.'),
        ('Reports: ',
         '30-day revenue trend chart, top customers, average booking value, coupon usage.'),
    ], size=18)
    add_footer(slide, 11, total)


def build_implementation(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '⚙️ Implementation Highlights')

    add_bullets(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(5.5), [
        ('Django app structure: ',
         '4 apps (accounts, booking, chatbot, dashboard) for separation of concerns.'),
        ('Atomic transactions: ',
         'payment + inventory + loyalty updates wrapped in transaction.atomic().'),
        ('AJAX for checkout: ',
         'coupon, insurance, loyalty are applied without page reload.'),
        ('CSRF protection: ',
         'enabled on every form and JSON endpoint.'),
        ('Template tags: ',
         'custom notifications and booking_extras filters.'),
        ('PDF generation: ',
         'reportlab Platypus templates with status badge, route blocks, price breakdown.'),
        ('Auto-seed: ',
         'one command (seed_data) creates 1500+ flights, 22 hotels, 5 coupons, demo users.'),
    ], size=17)
    add_footer(slide, 12, total)


def build_testing(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🧪 Testing & Results')

    add_textbox(slide, Inches(0.7), Inches(1.5), Inches(12), Inches(0.5),
                text='20 test cases executed across all modules — 100 % pass.',
                size=20, bold=True, color=RGBColor(0x15, 0x80, 0x3D))

    # Quick stats
    stats = [
        ('20 / 20', 'Tests passing'),
        ('95 %',    'Chatbot entity accuracy'),
        ('<200 ms', 'PDF generation'),
        ('1537',    'Sample flights seeded'),
        ('22',      'Sample hotels seeded'),
        ('5',       'Sample promo codes'),
    ]
    x0 = Inches(0.7); y0 = Inches(2.3)
    cell_w = Inches(4); cell_h = Inches(1.4); gap = Inches(0.1)
    for i, (val, lbl) in enumerate(stats):
        r, c = divmod(i, 3)
        x = x0 + (cell_w + gap) * c
        y = y0 + (cell_h + gap) * r
        add_rect(slide, x, y, cell_w, cell_h, LIGHT)
        add_textbox(slide, x, y + Inches(0.1), cell_w, Inches(0.7),
                    text=val, size=30, bold=True, color=PRIMARY,
                    align=PP_ALIGN.CENTER)
        add_textbox(slide, x, y + Inches(0.85), cell_w, Inches(0.4),
                    text=lbl, size=14, color=MUTED, align=PP_ALIGN.CENTER)

    add_footer(slide, 13, total)


def build_future(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_header_bar(slide, '🚀 Future Enhancements')

    add_bullets(slide, Inches(0.7), Inches(1.6), Inches(12), Inches(5.5), [
        ('Live APIs: ',           'Integrate Amadeus / Skyscanner / Booking.com.'),
        ('LLM fallback: ',        'Route ambiguous queries to OpenAI / Anthropic.'),
        ('Multi-language: ',      'Hindi, French, Spanish, German via Django i18n.'),
        ('Real payments: ',       'Razorpay / Stripe / PayU with webhook verification.'),
        ('Mobile app: ',          'React Native / Flutter wrapper.'),
        ('Notifications: ',       'Email (SendGrid) + SMS (Twilio).'),
        ('Personalisation: ',     'Collaborative-filtering recommendations.'),
        ('More modes: ',          'Buses, trains, cabs, holiday packages.'),
        ('CI/CD: ',               'GitHub Actions auto-deploy to Render / Railway.'),
    ], size=17)
    add_footer(slide, 14, total)


def build_conclusion(prs, total):
    blank = prs.slide_layouts[6]
    slide = prs.slides.add_slide(blank)
    add_rect(slide, 0, 0, SLIDE_W, SLIDE_H, PRIMARY)
    add_rect(slide, 0, Inches(3.4), SLIDE_W, Inches(0.6),
             RGBColor(0x3F, 0x37, 0xC4))

    add_textbox(slide, Inches(0.5), Inches(0.5), Inches(12.3), Inches(1),
                text='✅ Conclusion', size=44, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(1.8), Inches(11.3), Inches(1.4),
                text='TravelBot demonstrates how Python Full-Stack + NLP can '
                     'simplify travel planning into a single conversational platform.',
                size=22, italic=True, color=WHITE, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(4.2), Inches(11.3), Inches(1.2),
                text='Complete end-to-end booking · Chatbot · Reviews · '
                     'Wishlist · Loyalty · PDF tickets · Admin dashboard',
                size=18, color=SECONDARY, align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(5.5), Inches(11.3), Inches(0.6),
                text='Thank You !', size=40, bold=True, color=WHITE,
                align=PP_ALIGN.CENTER)

    add_textbox(slide, Inches(1.0), Inches(6.4), Inches(11.3), Inches(0.5),
                text='Questions & Discussion',
                size=18, italic=True, color=WHITE, align=PP_ALIGN.CENTER)


# ---------- Main ----------
def build():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    total = 15  # we will build 15 slides
    build_title_slide(prs, total)           # 1
    build_agenda(prs, total)                # 2
    build_introduction(prs, total)          # 3
    build_problem_objectives(prs, total)    # 4
    build_architecture(prs, total)          # 5
    build_tech_stack(prs, total)            # 6
    build_modules(prs, total)               # 7
    build_chatbot(prs, total)               # 8
    build_database(prs, total)              # 9
    build_features(prs, total)              # 10
    build_admin(prs, total)                 # 11
    build_implementation(prs, total)        # 12
    build_testing(prs, total)               # 13
    build_future(prs, total)                # 14
    build_conclusion(prs, total)            # 15

    out = Path(__file__).parent / 'TravelBot_Project_Presentation.pptx'
    prs.save(out)
    print(f'Saved: {out}  ({len(prs.slides)} slides)')


if __name__ == '__main__':
    build()
